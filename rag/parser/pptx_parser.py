"""
python-pptx를 사용한 PPTX 문서 텍스트 추출 모듈.

PPTX 파일의 각 슬라이드에서 텍스트 프레임(제목, 본문, 표 등)을 추출합니다.
슬라이드 번호와 함께 구조화된 텍스트를 반환합니다.
"""

import logging
from pathlib import Path

from rag.parser import EncryptedFileError

logger = logging.getLogger(__name__)


def _is_encrypted_ooxml(file_path: str) -> bool:
    """OOXML 파일이 MS-OFFCRYPTO 암호화 상태인지 확인 (OLE2 + EncryptedPackage)."""
    try:
        import olefile
        if olefile.isOleFile(file_path):
            ole = olefile.OleFileIO(file_path)
            encrypted = ole.exists("EncryptedPackage")
            ole.close()
            return encrypted
    except Exception:
        pass
    return False


def parse_pptx(file_path: str) -> str:
    """
    PPTX 파일의 모든 슬라이드에서 텍스트를 추출하여 하나의 문자열로 결합.

    각 슬라이드의 Shape에서 text_frame을 순회하며 텍스트를 수집합니다.
    테이블 Shape의 셀 텍스트도 추출합니다.

    Args:
        file_path: PPTX 파일 절대 경로.

    Returns:
        슬라이드별 텍스트를 빈 줄로 연결한 전체 텍스트 문자열.

    Raises:
        FileNotFoundError: 파일이 존재하지 않을 때.
        ImportError: python-pptx가 설치되지 않았을 때.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {file_path}")

    # 암호화 감지
    if _is_encrypted_ooxml(file_path):
        raise EncryptedFileError(file_path, "pptx")

    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Cannot parse PPTX files.")
        raise

    try:
        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            parts = [f"[슬라이드 {slide_num}]"]

            for shape in slide.shapes:
                # 텍스트 프레임이 있는 Shape (제목, 본문, 텍스트 박스 등)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)

                # 테이블 Shape
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            parts.append(" | ".join(row_texts))

            # 슬라이드에 텍스트가 있는 경우만 추가
            if len(parts) > 1:
                slides_text.append("\n".join(parts))

        full_text = "\n\n".join(slides_text)
        logger.info(
            f"PPTX parsed: {path.name} ({len(prs.slides)} slides, {len(full_text)} chars)"
        )
        return full_text

    except Exception:
        logger.exception(f"Failed to parse PPTX: {file_path}")
        raise
